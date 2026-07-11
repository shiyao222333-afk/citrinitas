"""
文本提取（按格式分发到各提取函数）
"""

import os
import re
import hashlib
import logging
import urllib.parse

logger = logging.getLogger(__name__)

from .registry import detect_file_type
from .encoding import _read_text_with_fallback


# ── 书籍内嵌图抽存 helper（#22：保留章节 + 抽出内嵌图，不 OCR，原图直存）──

def _safe_name(name: str) -> str:
    """把任意文件名清洗成安全文件名（去非法字符、截断过长）。"""
    name = re.sub(r'[\\/:*?"<>|]', "_", name).strip().strip(".")
    if not name:
        name = "image"
    return name[:120]


def _save_book_image(book_stem: str, name: str, data: bytes) -> str | None:
    """
    把书籍内嵌图原图直存到 library/images/books/<书名>/ 下。

    返回相对项目根的路径（如 library/images/books/LangChain/cover.jpg），
    供正文插入 [image:相对路径] 引用；失败返回 None（调用方跳过该图）。
    """
    try:
        from qconst import IMAGES_DIR, PROJECT_DIR
        sub = os.path.join(IMAGES_DIR, "books", _safe_name(book_stem))
        os.makedirs(sub, exist_ok=True)
        fname = _safe_name(name)
        dest = os.path.join(sub, fname)
        # 同名冲突但内容不同 → 加内容短哈希区分，避免覆盖
        if os.path.exists(dest) and hashlib.md5(data).hexdigest()[:8] != _file_md5(dest):
            base, ext = os.path.splitext(fname)
            dest = os.path.join(sub, f"{base}_{hashlib.md5(data).hexdigest()[:8]}{ext}")
        with open(dest, "wb") as f:
            f.write(data)
        # 统一正斜杠，保证跨平台一致且搜索卡片显示干净
        return os.path.relpath(dest, PROJECT_DIR).replace(os.sep, "/")
    except Exception as e:
        logger.warning(f"书籍图片保存失败（已跳过该图）: {e}")
        return None


def _file_md5(path: str) -> str:
    try:
        with open(path, "rb") as f:
            return hashlib.md5(f.read()).hexdigest()[:8]
    except Exception:
        return ""


def _norm_href(href: str) -> str:
    """EPUB 图片 href 归一化：去 #fragment + URL 解码，用于相对路径匹配。"""
    return urllib.parse.unquote(href.split("#")[0]).strip()


def _soup_to_structured_text(soup, book_stem: str, image_refs: list,
                             image_resolver=None) -> str:
    """
    把 BeautifulSoup 章节节点转成「保留标题层级 + 图片引用」的纯文本。

    - h1-h6 → Markdown 标题（# 到 ######）
    - img    → 抽原图存盘并插入 [image:相对路径]（不 OCR）
    - 其余块级元素 → 保留纯文本
    image_resolver: 可选回调(src, soup) -> bytes，用于 EPUB 把相对 src 解析为图片字节
    """
    out = []
    for elem in soup.find_all(["h1", "h2", "h3", "h4", "h5", "h6", "p", "img",
                               "figure", "li", "blockquote", "pre", "table"]):
        tag = elem.name
        if tag in ("h1", "h2", "h3", "h4", "h5", "h6"):
            level = int(tag[1])
            title = elem.get_text(strip=True)
            if title:
                out.append("#" * level + " " + title)
        elif tag == "img":
            src = elem.get("src") or elem.get("xlink:href")
            if src and not src.startswith("data:"):
                data = None
                if image_resolver:
                    try:
                        data = image_resolver(src, elem)
                    except Exception as e:
                        logger.warning(f"EPUB 图片解析失败 {src}: {e}")
                if data:
                    ref = _save_book_image(book_stem, os.path.basename(src) or "img", data)
                    if ref:
                        image_refs.append(ref)
                        out.append(f"[image:{ref}]")
        else:
            t = elem.get_text(separator="\n", strip=True)
            if t:
                out.append(t)
    return "\n".join(out)


def extract_text(file_path: str, file_info: dict = None) -> dict:
    """
    从文件提取纯文本内容。

    Args:
        file_path: 文件路径
        file_info: detect_file_type() 的结果，可选（None 时自动检测）

    Returns:
        {"ok": bool, "text": str, "error": str|None, "extraction_method": str}
    """
    if file_info is None:
        file_info = detect_file_type(file_path)
        if not file_info["ok"]:
            return {"ok": False, "text": "", "error": file_info["error"],
                    "extraction_method": None}

    fmt = file_info["format"]

    # 映射格式 → 提取函数
    extractors = {
        "epub": _extract_epub,
        "pdf":  _extract_pdf,
        "html": _extract_html,
        "htm":  _extract_html,
        "txt":  _extract_text_file,
        "md":   _extract_text_file,
        "json": _extract_text_file,
        "csv":  _extract_text_file,
        "srt":  _extract_srt,
        "docx": _extract_docx,
        "pptx": _extract_pptx,
    }

    extractor = extractors.get(fmt)
    if extractor is None:
        # 图片格式：返回空文本（留给 OCR 页面处理）
        if file_info["tier"] == 3:
            return {
                "ok": True,
                "text": "",
                "error": None,
                "extraction_method": "ocr_required",
                "ocr_required": True,
            }
        # 未知格式：兜底尝试
        try:
            text = _read_text_with_fallback(file_path)
            return {"ok": True, "text": text, "error": None,
                    "extraction_method": "direct_read_fallback"}
        except Exception as e:
            return {"ok": False, "text": "", "error": str(e), "extraction_method": None}

    try:
        result = extractor(file_path)
        if isinstance(result, str):
            result = {"ok": True, "text": result, "error": None}
        result.setdefault("extraction_method", file_info["extraction"])
        result.setdefault("ok", True)
        result.setdefault("error", None)
        return result
    except Exception as e:
        return {"ok": False, "text": "", "error": str(e),
                "extraction_method": file_info["extraction"]}


# ── 各格式提取函数 ──────────────────────────────────────────────────────────

def _extract_text_file(file_path: str) -> dict:
    """纯文本格式：TXT / MD / JSON / CSV"""
    text = _read_text_with_fallback(file_path)
    return {"ok": True, "text": text, "error": None}


def _extract_srt(file_path: str) -> dict:
    """SRT 字幕：去掉序号和时间戳，保留纯文本"""
    text = _read_text_with_fallback(file_path)
    # 去掉序号行和时间戳行
    lines = text.split("\n")
    cleaned = []
    for line in lines:
        line = line.strip()
        # 跳过序号（纯数字）
        if line.isdigit():
            continue
        # 跳过时间戳行 "00:00:01,000 --> 00:00:04,000"
        if re.match(r"^\d{2}:\d{2}:\d{2}[,.]\d{3}\s*-->\s*\d{2}:\d{2}:\d{2}[,.]\d{3}$", line):
            continue
        # 跳过 WEBVTT 头
        if line == "WEBVTT":
            continue
        if line:
            cleaned.append(line)
    return {"ok": True, "text": "\n".join(cleaned), "error": None}


def _extract_docx(file_path: str) -> dict:
    """Word 文档（.docx）：保留标题层级 + 抽出内嵌图（不 OCR，原图直存）"""
    try:
        import docx
    except ImportError:
        return {"ok": False, "text": "", "error": "请安装 python-docx: pip install python-docx"}
    try:
        doc = docx.Document(file_path)
    except Exception as e:
        return {"ok": False, "text": "", "error": f"DOCX 读取失败: {e}"}

    book_stem = os.path.splitext(os.path.basename(file_path))[0]
    lines = []
    image_refs = []

    # 正文：标题按级别转 Markdown #，其余段落保留
    for p in doc.paragraphs:
        text = p.text.strip()
        style = (p.style.name or "") if p.style else ""
        m = re.match(r"Heading\s+(\d+)", style)
        if m and text:
            lines.append("#" * int(m.group(1)) + " " + text)
        elif text:
            lines.append(text)

    # 内嵌图：统一抽到文末「插图」区（位置近似，不丢图）
    try:
        for i, shape in enumerate(doc.inline_shapes, 1):
            try:
                blip = shape._inline.graphic.graphicData.pic.blipFill.blip
                rId = blip.embed
                image_part = doc.part.related_parts[rId]
                ref = _save_book_image(book_stem, f"image_{i}.png", image_part.blob)
                if ref:
                    image_refs.append(ref)
            except Exception as e:
                logger.warning(f"DOCX 图片抽取失败（跳过）: {e}")
    except Exception as e:
        logger.warning(f"DOCX 图片遍历失败: {e}")

    if image_refs:
        lines.append("")
        lines.append("## 插图")
        lines.extend(f"[image:{r}]" for r in image_refs)

    text = "\n".join(lines)
    if not text.strip():
        return {"ok": False, "text": "", "error": "未能从 DOCX 中提取到文本内容"}
    return {"ok": True, "text": text, "error": None, "images": image_refs}


def _extract_pptx(file_path: str) -> dict:
    """PowerPoint 幻灯片（.pptx）"""
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        slides_text = []
        for i, slide in enumerate(prs.slides, 1):
            slide_lines = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        t = paragraph.text.strip()
                        if t:
                            slide_lines.append(t)
            if slide_lines:
                slides_text.append(f"--- 幻灯片 {i} ---\n" + "\n".join(slide_lines))
        text = "\n\n".join(slides_text)
        return {"ok": True, "text": text, "error": None}
    except ImportError:
        return {"ok": False, "text": "", "error": "请安装 python-pptx: pip install python-pptx"}


def _extract_html(file_path: str) -> dict:
    """HTML 网页：提取标题 + meta + body 文本"""
    text = _read_text_with_fallback(file_path)
    try:
        from bs4 import BeautifulSoup
        soup = BeautifulSoup(text, "html.parser")
        # 去掉 script/style
        for tag in soup(["script", "style", "nav", "footer", "header"]):
            tag.decompose()
        body_text = soup.get_text(separator="\n")
        # 清理多余空行
        lines = [line.strip() for line in body_text.split("\n") if line.strip()]
        return {"ok": True, "text": "\n".join(lines), "error": None}
    except ImportError:
        # 无 BeautifulSoup，返回原始 HTML 文本
        return {"ok": True, "text": text, "error": None,
                "warning": "未安装 BeautifulSoup，HTML 标签未清理"}


def _extract_epub(file_path: str) -> dict:
    """EPUB 电子书：保留章节结构（标题层级）+ 抽出内嵌图（不 OCR，原图直存）"""
    try:
        from bs4 import BeautifulSoup
    except ImportError:
        return {"ok": False, "text": "", "error": "请安装缺少的库: pip install beautifulsoup4"}
    from .epub_reader import read_epub, ITEM_IMAGE, ITEM_DOCUMENT

    try:
        book = read_epub(file_path)
    except Exception as e:
        return {"ok": False, "text": "", "error": f"EPUB 读取失败: {e}"}

    # 建图片 href 映射（归一化后用于相对路径解析）
    image_map = {}
    for it in book.get_items_of_type(ITEM_IMAGE):
        try:
            image_map[_norm_href(it.file_name)] = it
        except Exception:
            continue

    book_stem = os.path.splitext(os.path.basename(file_path))[0]
    chapters = []
    image_refs = []

    for item in book.get_items_of_type(ITEM_DOCUMENT):
        try:
            soup = BeautifulSoup(item.get_content(), "html.parser")
            for tag in soup(["script", "style"]):
                tag.decompose()
            base_href = item.file_name

            # 把相对 src 解析为图片字节
            def _resolver(src, _elem, _base=base_href, _map=image_map):
                rel = urllib.parse.urljoin(_base, src)
                target = _map.get(_norm_href(rel))
                if target is None:
                    return None
                return target.get_content()

            block = _soup_to_structured_text(soup, book_stem, image_refs, _resolver)
            if block.strip():
                chapters.append(block)
        except Exception:
            continue

    if not chapters:
        return {"ok": False, "text": "",
                "error": "未能从 EPUB 中提取到文本内容"}
    text = "\n\n".join(chapters)
    return {"ok": True, "text": text, "error": None, "images": image_refs}


def _extract_pdf(file_path: str) -> dict:
    """
    PDF 双路径提取:
    1. 先用 pypdf 提取文字层
    2. 文字层不足 → 标记为需要 OCR
    """
    try:
        from pypdf import PdfReader
    except ImportError:
        # pypdf 可能叫 PyPDF2
        try:
            from PyPDF2 import PdfReader
        except ImportError:
            return {"ok": False, "text": "",
                    "error": "请安装 pypdf: pip install pypdf"}

    try:
        reader = PdfReader(file_path)
    except Exception as e:
        return {"ok": False, "text": "", "error": f"PDF 文件无法打开: {e}"}

    # 提取所有页的文本 + 内嵌图（原图直存，不 OCR）
    book_stem = os.path.splitext(os.path.basename(file_path))[0]
    pages_text = []
    image_refs = []
    for i, page in enumerate(reader.pages):
        try:
            pt = page.extract_text() or ""
        except Exception:
            pt = ""
        # 抽该页图片
        try:
            for img in page.images:
                data = getattr(img, "data", None)
                if data:
                    name = getattr(img, "name", f"page{i+1}_img") or f"page{i+1}_img"
                    ref = _save_book_image(book_stem, name, data)
                    if ref:
                        image_refs.append(ref)
                        pt = (pt + f"\n\n[image:{ref}]").strip()
        except Exception as e:
            logger.warning(f"PDF 第{i+1}页图片抽取失败（跳过）: {e}")
        if pt and pt.strip():
            pages_text.append(f"--- 第 {i+1} 页 ---\n" + pt.strip())

    text = "\n\n".join(pages_text)

    # 判断是否有足够的文字层
    # 规则：总字符数 < 100 且页数 > 1 → 可能是扫描版
    if len(text) < 100 and len(reader.pages) > 1:
        return {
            "ok": True,
            "text": text or "",
            "error": None,
            "images": image_refs,
            "warning": "此 PDF 文字层很少，可能是扫描版，建议使用 OCR 识别。",
            "ocr_recommended": True,
            "total_pages": len(reader.pages),
            "has_text_layer": False,
        }

    return {
        "ok": True,
        "text": text,
        "error": None,
        "images": image_refs,
        "total_pages": len(reader.pages),
        "has_text_layer": bool(text.strip()),
    }
